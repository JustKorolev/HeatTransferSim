"""Experimental-setup slide: an editable .pptx table plus a preview PNG.

Two outputs from one definition of the rows, because a hand-kept copy of the same
table in two formats is a copy that goes stale -- and several of these numbers are
measured values that have already been wrong once.

    experimental_setup.pptx          native table, editable text, black slide
    experimental_setup_preview.png   the same thing rendered, for checking

Values still to be filled from the run's own simulation_parameters.json are
written as ``___`` and rendered in amber so they cannot be presented by accident.

    python docs/surf_report/make_setup_table.py
"""

from __future__ import annotations

from pathlib import Path

# ---------------------------------------------------------------- the content
# (label, value). "___" marks a value nobody has supplied yet; it renders amber.
SECTIONS: list[tuple[str, list[tuple[str, str]]]] = [
    ("Plant & model", [
        ("Graph", "2.99 M cells · 27 heaters · 27 sensors"),
        ("Gain G", "exact DC solve of  L T = P"),
        ("Linearised at", "T_op = 50 K"),
        ("Conditioning", "cond 366 · σ₁ = 81 % of the plant"),
        ("Properties", "k(T), c_p(T)  ___"),
        ("Radiation", "to ambient only — no view factors"),
        ("Solver", "implicit step, GPU conjugate gradient"),
        ("Conservation", "energy drift ~1e−6 throughout"),
    ]),
    ("Controller & run", [
        ("Scheme", "static-decoupling MIMO PI"),
        ("Allocation", "bounded least squares,  u ≥ 0"),
        ("Gains", "K_p = 5.5 · K_i = ___"),
        ("Filter", "τ_f = 900 s, feedback path only"),
        ("Step size", "dt = ___ s"),
        ("Heater limit", "___ W each"),
        ("Setpoints", "___ K, uniform across sensors"),
        ("Channels", "25 of 27 in the loop — 2 unreachable"),
        ("Initial state", "uniform at ___ K"),
        ("Duration", "27.8 h ≈ 1.2 × dominant time constant"),
    ]),
]

TITLE = "Experimental setup"
SUBTITLE = "Cryostat thermal simulator — what was held fixed, and at what value"
NOTES = (
    "Amber entries are still placeholders: pull dt, K_i, the heater limit, the setpoint "
    "and the initial temperature from the run's simulation_parameters.json before "
    "presenting. The heater limit especially — the saved parameter says 1.5 W per heater "
    "while the talk quotes 7-11 W of 30, and those cannot both be right."
)

# ------------------------------------------------------------------- palette
INK = "F5F5F5"
LABEL = "B8B8B8"
TODO = "F0A95C"
HEADER_FILL = "142C3C"
HEADER_INK = "FFFFFF"
ROW_A, ROW_B = "0C0C0C", "151515"
BORDER = "333333"
SUB = "9A9A9A"
BG = "000000"

TITLE_FONT, BODY_FONT = "Cambria", "Calibri"

# Shared by both builders so the .pptx and its preview cannot drift. Sized to fill
# the slide rather than to fit: this is read from the back of a room, and the
# first pass left the bottom third empty at font sizes nobody could see.
HEADER_H_IN, ROW_H_IN = 0.50, 0.44
LABEL_W_IN, VALUE_W_IN = 1.95, 4.05
TABLE_X_IN, TABLE_Y_IN, COLUMN_STEP_IN = 0.55, 1.62, 6.34
HEADER_PT, LABEL_PT, VALUE_PT = 16.0, 13.0, 14.0


def is_todo(value: str) -> bool:
    return "___" in value


# ----------------------------------------------------------------------- pptx
def build_pptx(target: Path) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    def rgb(value: str) -> RGBColor:
        return RGBColor.from_string(value)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])       # blank

    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(BG)

    def text_box(x, y, w, h, text, size, color, font, bold=False):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
        run = frame.paragraphs[0].add_run()
        run.text = text
        run.font.size, run.font.bold = Pt(size), bold
        run.font.color.rgb, run.font.name = rgb(color), font

    text_box(0.55, 0.42, 12.2, 0.7, TITLE, 34, INK, TITLE_FONT, bold=True)
    text_box(0.55, 1.06, 12.2, 0.4, SUBTITLE, 14, SUB, BODY_FONT)

    for column, (heading, rows) in enumerate(SECTIONS):
        left = TABLE_X_IN + column * COLUMN_STEP_IN
        table_shape = slide.shapes.add_table(
            len(rows) + 1, 2,
            Inches(left), Inches(TABLE_Y_IN), Inches(LABEL_W_IN + VALUE_W_IN),
            Inches(HEADER_H_IN + ROW_H_IN * len(rows)),
        )
        table = table_shape.table
        # python-pptx applies a banded style by default; every cell is filled
        # explicitly below so the theme cannot reintroduce light rows on a black
        # slide.
        table.first_row = False
        table.horz_banding = False
        table.columns[0].width = Inches(LABEL_W_IN)
        table.columns[1].width = Inches(VALUE_W_IN)
        table.rows[0].height = Inches(HEADER_H_IN)

        for r in range(len(rows) + 1):
            if r:
                table.rows[r].height = Inches(ROW_H_IN)
            for c in range(2):
                cell = table.cell(r, c)
                cell.margin_left = Inches(0.11)
                cell.margin_right = Inches(0.08)
                cell.margin_top = cell.margin_bottom = Inches(0.02)
                fill = cell.fill
                fill.solid()
                if r == 0:
                    fill.fore_color.rgb = rgb(HEADER_FILL)
                else:
                    fill.fore_color.rgb = rgb(ROW_A if r % 2 else ROW_B)

                if r == 0:
                    text, size, color, bold = (heading if c == 0 else ""), HEADER_PT, HEADER_INK, True
                else:
                    label, value = rows[r - 1]
                    if c == 0:
                        text, size, color, bold = label, LABEL_PT, LABEL, False
                    else:
                        text, size, bold = value, VALUE_PT, False
                        color = TODO if is_todo(value) else INK

                frame = cell.text_frame
                frame.word_wrap = True
                paragraph = frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = text
                run.font.size, run.font.bold = Pt(size), bold
                run.font.color.rgb, run.font.name = rgb(color), BODY_FONT

    slide.notes_slide.notes_text_frame.text = NOTES
    prs.save(target)
    return target


# ------------------------------------------------------------------ preview
def build_preview(target: Path) -> Path:
    """Render the same table with matplotlib, purely so the layout is checkable.

    python-pptx cannot rasterise, so without this the .pptx ships unseen -- and a
    table that overflows its column is invisible in the XML.
    """
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import Rectangle

    W, H = 1333.0, 750.0                      # 13.33 x 7.5 in, at 100 px/in
    fig = plt.figure(figsize=(W / 100.0, H / 100.0))
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W)
    ax.set_ylim(H, 0)
    ax.axis("off")
    ax.add_patch(Rectangle((0, 0), W, H, facecolor="#" + BG, zorder=0))

    ax.text(55, 62, TITLE, fontsize=25, fontweight="bold", color="#" + INK,
            va="center", family="serif", zorder=2)
    ax.text(55, 116, SUBTITLE, fontsize=10.5, color="#" + SUB, va="center", zorder=2)

    # 100 px per inch, so the shared inch-based layout carries over unchanged and
    # a point is 100/72 px.
    px, pt = 100.0, 100.0 / 72.0
    for column, (heading, rows) in enumerate(SECTIONS):
        left = TABLE_X_IN * px + column * COLUMN_STEP_IN * px
        label_w, value_w = LABEL_W_IN * px, VALUE_W_IN * px
        y = TABLE_Y_IN * px
        header_h, row_h = HEADER_H_IN * px, ROW_H_IN * px
        ax.add_patch(Rectangle((left, y), label_w + value_w, header_h,
                               facecolor="#" + HEADER_FILL, edgecolor="#" + BORDER,
                               linewidth=0.8, zorder=1))
        ax.text(left + 11, y + header_h / 2, heading, fontsize=HEADER_PT * pt * 0.72,
                fontweight="bold", color="#" + HEADER_INK, va="center", zorder=2)
        y += header_h
        for i, (label, value) in enumerate(rows):
            band = ROW_A if (i + 1) % 2 else ROW_B
            ax.add_patch(Rectangle((left, y), label_w + value_w, row_h,
                                   facecolor="#" + band, edgecolor="#" + BORDER,
                                   linewidth=0.8, zorder=1))
            ax.text(left + 11, y + row_h / 2, label, fontsize=LABEL_PT * pt * 0.72,
                    color="#" + LABEL, va="center", zorder=2)
            ax.text(left + label_w + 11, y + row_h / 2, value, fontsize=VALUE_PT * pt * 0.72,
                    color="#" + (TODO if is_todo(value) else INK), va="center", zorder=2)
            y += row_h

    fig.savefig(target, dpi=200, facecolor="#" + BG)
    plt.close(fig)
    return target


def main() -> None:
    here = Path(__file__).resolve().parent
    print(f"wrote {build_pptx(here / 'experimental_setup.pptx')}")
    print(f"wrote {build_preview(here / 'experimental_setup_preview.png')}")


if __name__ == "__main__":
    main()
